#!/usr/bin/env python3
"""
File type extractors for RAG pipeline - Jina v5 pure text version.

- Image files -> PaddleOCR CPU -> text -> embedding
- Image-only PDFs -> PaddleOCR CPU -> text -> embedding
- OCR quality check: garbled/too-short results are skipped
- Video files: SKIPPED (no longer processed)
- All chunks are text (no image/video embedding)

Configuration loaded from config.yaml via config_loader.
"""
import os
import sys
import re
import json
import tempfile
import logging
import sqlite3
import threading
import shutil
import subprocess
import multiprocessing
import queue
from typing import List, Optional, Tuple
from dataclasses import dataclass, field

# Import config loader for centralized configuration
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
try:
    from config_loader import (
        get_db_path, get_log_dir, get_ocr_api_url,
        get_chunking_config, get_limits_config, get_config
    )
    config = get_config()
    DB_PATH = get_db_path()
    LOG_DIR = get_log_dir()
    _OCR_API_URL = get_ocr_api_url()
    MIN_TEXT_CHARS = config['chunking']['min_text_chars']
    MIN_TEXT_CHARS_LOW = config['chunking']['min_text_chars_low']
    limits = get_limits_config()
    MIN_MEDIA_BYTES = limits['min_media_bytes']
    MAX_IMAGE_BYTES = limits['max_image_bytes']
    MAX_TEXT_FILE_BYTES = limits['max_text_file_bytes']
    MAX_DJVU_PAGES = limits['max_djvu_pages']
except ImportError:
    # Fallback for standalone usage
    DB_PATH = os.environ.get("DB_PATH", "./file_index.db")
    _OCR_API_URL = os.environ.get("OCR_API_URL", "http://127.0.0.1:8002")
    MIN_TEXT_CHARS = 500        # Standard minimum: Skip extraction below 500 chars
    MIN_TEXT_CHARS_LOW = 50     # Lower threshold for images/screenshots/logos/business cards
    MIN_MEDIA_BYTES = 10 * 1024  # Pre-filter: Skip images/audio smaller than 10KB
    MAX_IMAGE_BYTES = 100 * 1024 * 1024  # Pre-filter: Skip images larger than 100MB
    MAX_TEXT_FILE_BYTES = 50 * 1024 * 1024  # 50MB，超过则跳过
    MAX_DJVU_PAGES = 500

from chunker_v2 import chunk_greedy_semantic


log = logging.getLogger(__name__)

# --- PaddleOCR: subprocess isolation ---
# 每次调用 PaddleOCR 都在独立子进程中运行（CUDA_VISIBLE_DEVICES='' 禁用 CUDA）
# 子进程退出后内存完全释放，从根本上消除 PaddlePaddle GPU 版本的内存泄漏
# PaddlePaddle GPU 版即使 use_gpu=False 也会加载 CUDA runtime + 预分配 64GB VA

os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"

_OCR_LOCK = threading.RLock()  # 同一进程内线程安全


def _run_ocr_on_images(image_paths: List[str]) -> List[Optional[str]]:
    """Run PaddleOCR on a list of images via OCR API server.
    
    Calls the persistent OCR API server (ocr_server.py) which loads
    PaddleOCR once into GPU memory and reuses it across calls.
    Returns text for each image (None if failed).
    """
    if not image_paths:
        return []
    
    try:
        with _OCR_LOCK:
            import requests as _req
            resp = _req.post(
                f"{_OCR_API_URL}/ocr",
                json={"images": image_paths},
                timeout=600,
                proxies={"http": None, "https": None},
            )
        
        if resp.status_code != 200:
            log.warning(f"OCR API error (status={resp.status_code}): {resp.text[:300]}")
            return [None] * len(image_paths)
        
        data = resp.json()
        results = data.get('results', [])
        
        if len(results) != len(image_paths):
            log.warning(f"OCR result count mismatch: got {len(results)}, expected {len(image_paths)}")
            return [None] * len(image_paths)
        
        return results
        
    except Exception as e:
        log.warning(f"OCR API call failed: {e}")
        return [None] * len(image_paths)


def _ocr_image(image_path: str) -> Optional[str]:
    """Run OCR on a single image via subprocess.
    
    Returns extracted text, or None if OCR fails / quality too low.
    """
    if not os.path.isfile(image_path):
        return None
    
    results = _run_ocr_on_images([image_path])
    
    if not results:
        return None
    return results[0]


def _ocr_images_batch(image_paths: List[str]) -> List[Optional[str]]:
    """Run OCR on multiple images via subprocess.
    
    Batching: max 20 images per subprocess call (PaddleOCR ~3-5s/image, 
    subprocess timeout 300s → 20 images max to stay safe).
    Returns list of extracted texts (None for failed/too-low-quality).
    """
    if not image_paths:
        return []

    all_results = []
    batch_size = 20
    total_batches = (len(image_paths) + batch_size - 1) // batch_size
    
    for batch_idx in range(total_batches):
        batch_start = batch_idx * batch_size
        batch_end = min(batch_start + batch_size, len(image_paths))
        batch_paths = image_paths[batch_start:batch_end]
        
        log.info(f"OCR batch {batch_idx+1}/{total_batches}: {len(batch_paths)} images")
        results = _run_ocr_on_images(batch_paths)
        all_results.extend(results)
    
    return all_results


def _is_ocr_garbage(text: str, min_chars: int = 10) -> bool:
    """Check if OCR result is garbled or too short to be useful.

    Criteria:
    - Too short (< min_chars)
    - High ratio of non-alphanumeric characters (garbage/encoding errors)
    - Repetitive characters (OCR artifacts)
    - No meaningful words detected

    Returns True if text should be skipped (is garbage).
    """
    if not text or len(text.strip()) < min_chars:
        return True

    text = text.strip()

    # Count meaningful characters
    alnum_count = sum(1 for c in text if c.isalnum())
    total_count = len(text.replace('\n', '').replace(' ', ''))

    if total_count == 0:
        return True

    alnum_ratio = alnum_count / total_count

    # If less than 30% alphanumeric, likely garbage
    if alnum_ratio < 0.3:
        return True

    # Check for repetitive OCR artifacts (same char repeated 10+ times)
    # Allow legitimate patterns like separators, table fillers
    # 用正则优化，避免 O(n) 循环
    if re.search(r'([a-zA-Z\u4e00-\u9fff])\1{9,}', text):
        # 中文/字母重复10次以上才算垃圾(允许分隔线如 '------')
        return True

    # Check for excessive unique characters (random byte noise)
    unique_chars = len(set(text.replace('\n', '').replace(' ', '')))
    if total_count > 50 and unique_chars / total_count > 0.90:
        # Almost every char is unique - likely noise
        return True

    return False


# --- Whisper: process-global model (thread-safe with lock) ---
# Changed from per-thread to process-global to avoid repeated loading
_WHISPER_MODEL = None
_WHISPER_LOCK = threading.Lock()


def _get_whisper():
    """Lazy-load Whisper model (process-global, thread-safe with lock)."""
    global _WHISPER_MODEL
    if _WHISPER_MODEL is None:
        with _WHISPER_LOCK:
            if _WHISPER_MODEL is None:  # double-check locking
                from faster_whisper import WhisperModel
                _WHISPER_MODEL = WhisperModel("small", device="cpu", compute_type="int8")
                log.info(f"✅ Whisper model loaded (process-global, cached for reuse)")
    return _WHISPER_MODEL


# --- whisper.cpp subprocess (faster fallback if available) ---
_WHISPER_CPP_PATH = None

def _init_whisper_cpp():
    """Find whisper-cli binary. Returns path or None."""
    global _WHISPER_CPP_PATH
    if _WHISPER_CPP_PATH is not None:
        return _WHISPER_CPP_PATH
    for candidate in [
        shutil.which('whisper-cli'),
        shutil.which('main'),
        os.path.expanduser('~/whisper.cpp/build/bin/main'),
        os.path.expanduser('~/whisper.cpp/build/bin/whisper-cli'),
    ]:
        if candidate and os.path.isfile(candidate):
            _WHISPER_CPP_PATH = candidate
            return candidate
    _WHISPER_CPP_PATH = False
    return None


def _transcribe_with_whisper_cpp(audio_path: str) -> Optional[Tuple[str, str]]:
    """Transcribe audio using whisper.cpp subprocess."""
    cpp_path = _init_whisper_cpp()
    if not cpp_path:
        return None
    model_path = os.path.expanduser("~/models/ggml-small.bin")
    if not os.path.isfile(model_path):
        model_path = os.path.expanduser("~/models/ggml-small.en.bin")
    if not os.path.isfile(model_path):
        return None
    try:
        r = subprocess.run(
            [cpp_path, '-m', model_path, '-f', audio_path, '-otxt', '-l', 'auto'],
            capture_output=True, text=True, timeout=300
        )
        if r.returncode == 0:
            txt_path = audio_path + '.txt'
            if os.path.isfile(txt_path):
                with open(txt_path, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read().strip()
                try:
                    os.remove(txt_path)
                except OSError:
                    pass
                if text:
                    return (text, 'unknown')
    except (subprocess.TimeoutExpired, OSError, Exception) as e:
        log.debug(f"whisper.cpp failed for {audio_path}: {e}")
    return None


# --- Music-only detection heuristic ---
_MUSIC_PATH_HINTS = {
    '/music/', '/music library', '/songs/', '/tracks/', '/album/',
    '/soundtrack/', '/ost/', '/sound effects/', '/sfx/',
    '/音频/', '/音乐/', '/歌曲/', '/原声/', '/音效/',
}
_MUSIC_EXT_HINTS = {'mp3', 'flac', 'wav', 'ogg', 'opus', 'm4a'}
_MUSIC_DIR_HINTS = {'music', 'songs', 'tracks', 'albums', 'ost', 'soundtrack', 'audio'}


def _is_likely_music_only(path: str, file_size: int = 0, media_meta: str = None) -> bool:
    """Heuristic: skip Whisper for paths that look like pure music/SFX."""
    path_lower = path.lower()
    for hint in _MUSIC_PATH_HINTS:
        if hint in path_lower:
            return True
    ext = path.rsplit('.', 1)[-1].lower() if '.' in path else ''
    if ext in _MUSIC_EXT_HINTS:
        parts = path_lower.split('/')
        for part in parts:
            if part in _MUSIC_DIR_HINTS:
                return True
    if media_meta:
        try:
            meta = json.loads(media_meta) if isinstance(media_meta, str) else media_meta
            if isinstance(meta, dict):
                tags = meta.get('tags', {})
                if isinstance(tags, dict):
                    genre = str(tags.get('genre', '')).lower()
                    if genre in ('music', 'instrumental', 'soundtrack', 'ost'):
                        return True
        except (json.JSONDecodeError, TypeError, AttributeError):
            pass
    return False


# --- Data Types ---
@dataclass
class ChunkData:
    """A single chunk ready for embedding.
    
    Jina v5 pipeline: all chunks are text (OCR/Whisper converts media to text).
    """
    file_id: int
    chunk_index: int
    text: str
    char_count: int = 0
    meta: dict = field(default_factory=dict)  # {'source': 'pdf_ocr', 'ocr_pages': 5, 'language': 'zh'}


# --- File Type Classification ---
TEXT_EXTS = {'.txt', '.md', '.rst', '.log', '.csv', '.tsv', '.lrc'}
CODE_EXTS = {'.py', '.js', '.java', '.c', '.cpp', '.h', '.hpp', '.cs', '.go',
             '.rs', '.rb', '.php', '.lua', '.sh', '.sql', '.r', '.m', '.scala', '.kt',
             '.swift', '.dart', '.ps1', '.bat', '.css', '.scss', '.html', '.vue', '.jsx', '.tsx'}
CONFIG_EXTS = {'.json', '.yaml', '.yml', '.toml', '.xml', '.ini', '.cfg', '.conf',
               '.env', '.properties', '.gitignore', '.dockerignore', '.editorconfig'}
PDF_EXTS = {'.pdf'}
IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif', '.svg', '.ico'}
VIDEO_EXTS = {'.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv', '.webm', '.m4v', '.3gp', '.ts'}
AUDIO_EXTS = {'.mp3', '.wav', '.flac', '.ogg', '.m4a', '.aac', '.wma', '.opus', '.wem', '.amr', '.aiff', '.aif', '.sfk', '.mid', '.midi', '.rmi'}
# Binary/garbage extensions: skip entirely (never try to read as text)
BINARY_EXTS = {
    '.bin', '.exe', '.elf', '.dll', '.so', '.dylib', '.o', '.obj', '.pyc', '.pyo',
    '.dat', '.tar.enc', '.dds', '.tga', '.bundle', '.ocl', '.majset', '.psobin',
    '.slp', '.dxs', '.w3x', '.idx', '.fxo', '.blob', '.shot', '.thumb',
    '.null', '.ds_store', '.3dmark-result', '.fastresume',
    '.gz.cdx', '.gz.krt',
}
# Numbered extensions (e.g. .0, .1, .2 ... .13) — game cache / partition data
BINARY_EXTS.update(f'.{i}' for i in range(20))

EBOOK_EXTS = {'.epub', '.mobi', '.azw3', '.fb2'}
DJVU_EXTS = {'.djvu'}  # djvu 跟 PDF 相同逻辑：文本 → OCR → 视觉模型
OFFICE_EXTS = {'.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.odt', '.ods', '.odp', '.rtf'}
ARCHIVE_EXTS = {'.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.xz'}


def classify_file(ext: str, media_meta: Optional[str] = None) -> str:
    """Classify file into processing type."""
    ext = ext.lower()
    if ext in TEXT_EXTS: return 'text'
    if ext in CODE_EXTS: return 'code'
    if ext in CONFIG_EXTS: return 'config'
    if ext in IMAGE_EXTS: return 'image'
    if ext in VIDEO_EXTS: return 'video'
    if ext in AUDIO_EXTS: return 'audio'
    if ext in EBOOK_EXTS: return 'ebook'
    if ext in DJVU_EXTS: return 'djvu'
    if ext in OFFICE_EXTS: return 'office'
    if ext in ARCHIVE_EXTS: return 'archive'
    if ext in BINARY_EXTS: return 'binary'  # 放在最后，让特定类型优先匹配
    if ext == '.pdf':
        if media_meta:
            try:
                meta = json.loads(media_meta) if isinstance(media_meta, str) else media_meta
                if isinstance(meta, dict):
                    ptype = meta.get('type', '')
                    if ptype == 'text_pdf':
                        return 'pdf_text'
            except (json.JSONDecodeError, TypeError, AttributeError):
                pass
        return 'pdf_image'
    return 'unknown'


# File types that should use lower OCR threshold (scanned images, screenshots, logos, business cards)
LOW_OCR_THRESHOLD_EXTS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif', '.ico',
                          '.pdf'}  # Also lower threshold for PDFs (charts, forms, scanned docs)


def _mark_ocr_insufficient(file_id: int, path: str, reason: str):
    """Mark a file in the DB as OCR-insufficient so it's not re-processed."""
    try:
        conn = sqlite3.connect(DB_PATH, timeout=30)
        c = conn.cursor()
        # Create status table if not exists
        c.execute("""CREATE TABLE IF NOT EXISTS rag_ocr_status (
            file_id INTEGER PRIMARY KEY,
            status TEXT,  -- 'insufficient', 'garbage', 'timeout', 'error'
            reason TEXT,
            checked_at TEXT DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (file_id) REFERENCES files(id)
        )""")
        c.execute("INSERT OR REPLACE INTO rag_ocr_status (file_id, status, reason) VALUES (?, ?, ?)",
                  (file_id, 'insufficient', reason))
        conn.commit()
        conn.close()
        log.info(f"🏷️ Marked OCR-insufficient: file_id={file_id} reason={reason} path={path}")
    except Exception as e:
        log.warning(f"Failed to mark OCR-insufficient for file_id={file_id}: {e}")


def extract(file_id: int, file_path: str, ext: str, file_type: str,
            media_meta: Optional[str] = None, file_size: int = 0) -> List[ChunkData]:
    """Main extraction dispatcher."""
    if not os.path.exists(file_path):
        log.warning(f"File not found: {file_path}")
        # 标记为已删除，后续运行跳过
        try:
            conn = sqlite3.connect(DB_PATH, timeout=5)
            conn.execute("UPDATE files SET is_deleted=1 WHERE id=?", (file_id,))
            conn.commit()
            conn.close()
        except Exception:
            pass
        return []

    # Pre-filter: Skip tiny images/audio (save GPU resources)
    if file_type in ('image', 'audio') and file_size < MIN_MEDIA_BYTES:
        return []

    # Pre-filter: Skip huge images (avoid OCR timeout/GPU memory)
    # 100MB threshold: larger images cause PaddleOCR to timeout or crash
    if file_type == 'image' and file_size > MAX_IMAGE_BYTES:
        log.debug(f"Image too large ({file_size//1024//1024}MB > {MAX_IMAGE_BYTES//1024//1024}MB): {file_path}")
        return []

    # Skip video files - not processed in Jina v5 pipeline
    if file_type == 'video':
        return []

    try:
        if file_type in ('text', 'code', 'config'):
            return _extract_text(file_id, file_path, ext, media_meta, file_type=file_type)
        elif file_type == 'pdf_text':
            return _extract_pdf_text(file_id, file_path, media_meta)
        elif file_type == 'pdf_image':
            return _extract_pdf_image_ocr(file_id, file_path, media_meta)
        elif file_type == 'image':
            return _extract_image_ocr(file_id, file_path, media_meta)
        elif file_type == 'djvu':
            return _extract_djvu(file_id, file_path, media_meta)
        elif file_type == 'video':
            return []  # Skip videos entirely
        elif file_type == 'audio':
            return _extract_audio(file_id, file_path, file_size, media_meta)
        elif file_type == 'ebook':
            return _extract_ebook(file_id, file_path, ext, media_meta)
        elif file_type == 'office':
            return _extract_office(file_id, file_path, ext, media_meta)
        elif file_type == 'archive':
            return _extract_archive(file_id, file_path, media_meta)
        elif file_type == 'binary':
            return []  # Skip binary files (already filtered by classify_file)
        else:
            return _extract_text_fallback(file_id, file_path, media_meta)
    except Exception as e:
        log.warning(f"Extract error for {file_path}: {e}")
        return []


# ============================================================
# Individual Extractors
# ============================================================

def _build_context_text(path: str, media_meta: Optional[str] = None, short: bool = False) -> str:
    """Build context text: path + filename + metadata."""
    parts = []

    if short:
        path_parts = path.rsplit('/', 3)
        short_path = '/' + '/'.join(path_parts[-3:]) if len(path_parts) > 3 else path
        parts.append(f"File: ...{short_path}")
    else:
        parts.append(f"File: {path}")

    if media_meta:
        try:
            meta = json.loads(media_meta) if isinstance(media_meta, str) else media_meta
            if isinstance(meta, dict):
                meta_lines = []
                for k, v in meta.items():
                    if isinstance(v, dict):
                        for k2, v2 in v.items():
                            meta_lines.append(f"  {k}.{k2}: {v2}")
                    elif v is not None:
                        meta_lines.append(f"  {k}: {v}")
                if meta_lines:
                    parts.append("Metadata:\n" + "\n".join(meta_lines))
        except (json.JSONDecodeError, TypeError, AttributeError):
            pass

    return "\n".join(parts)


def _read_file(path: str, max_chars: int = 0) -> Optional[str]:
    """Read text file with encoding detection. Skips files > MAX_TEXT_FILE_BYTES to avoid OOM."""
    try:
        file_size = os.path.getsize(path)
    except OSError:
        return None

    # Hard skip: files > MAX_TEXT_FILE_BYTES are likely binary/cache/game data, not text
    if file_size > MAX_TEXT_FILE_BYTES:
        log.debug(f"Skipped large file ({file_size // 1024 // 1024}MB > {MAX_TEXT_FILE_BYTES // 1024 // 1024}MB): {path}")
        return None

    cap = max_chars if max_chars > 0 else None
    for enc in ['utf-8', 'gb18030', 'latin-1']:
        try:
            with open(path, 'r', encoding=enc, errors='replace') as f:
                if cap:
                    text = f.read(cap)
                else:
                    text = f.read()
            # 过滤 null bytes（UTF-16LE 误读为 UTF-8 时会出现大量 \x00）
            if text and '\x00' in text:
                null_count = text.count('\x00')
                if null_count > 50:  # 超过 50 个 null bytes = 编码错误/二进制
                    log.debug(f"Skipped null-byte file ({null_count} nulls): {path}")
                    return None
                text = text.replace('\x00', '')  # 少量直接清除
            return text
        except (UnicodeDecodeError, UnicodeError):
            continue
        except OSError:
            return None
    return None


def _extract_text(file_id: int, path: str, ext: str, media_meta: str = None,
                  file_type: str = 'text') -> List[ChunkData]:
    """Text/code/config files: read -> greedy semantic chunk with context prefix."""
    text = _read_file(path)
    if not text or len(text.strip()) < MIN_TEXT_CHARS:
        _mark_ocr_insufficient(file_id, path, f'text_short_{len(text.strip()) if text else 0}chars')
        return []

    # Garbled text check (encoding errors, binary garbage)
    if _is_ocr_garbage(text.strip(), min_chars=MIN_TEXT_CHARS):
        _mark_ocr_insufficient(file_id, path, 'text_garbled')
        return []

    ctx = _build_context_text(path, media_meta)
    full_text = ctx + "\n\n" + text

    chunks = chunk_greedy_semantic(full_text, ext=ext)


    return [ChunkData(file_id=file_id, chunk_index=i, text=c, char_count=len(c)) for i, c in enumerate(chunks)]


def _extract_pdf_text(file_id: int, path: str, media_meta: Optional[str]) -> List[ChunkData]:
    """Text PDF: PyMuPDF extract text -> greedy chunk."""
    try:
        import fitz
    except ImportError:
        log.warning("PyMuPDF not installed, skipping PDF text extraction")
        return []

    # PyMuPDF 不支持非 ASCII 路径，用 stream 模式打开
    try:
        with open(path, 'rb') as f:
            doc = fitz.open(stream=f.read(), filetype='pdf')
    except Exception as e:
        log.warning(f"Failed to open PDF {path}: {e}")
        return []
    
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    try:
        doc.close()
    except Exception as e:
        log.warning(f"Error closing PDF: {e}")

    if len(text.strip()) < MIN_TEXT_CHARS:
        _mark_ocr_insufficient(file_id, path, f'pdf_text_short_{len(text.strip())}chars')
        return []

    # Garbled text check
    if _is_ocr_garbage(text.strip(), min_chars=MIN_TEXT_CHARS):
        _mark_ocr_insufficient(file_id, path, 'pdf_text_garbled')
        return []

    chunks = chunk_greedy_semantic(text, ext='.txt')
    result = []
    short_ctx = _build_context_text(path, media_meta, short=True)
    for i, chunk_text in enumerate(chunks):
        full_chunk = f"{short_ctx}\n\n{chunk_text}"
        result.append(ChunkData(
            file_id=file_id, chunk_index=i,
            text=full_chunk, char_count=len(full_chunk),
            meta={'source': 'pdf_text'}
        ))
    return result


def _extract_pdf_image_ocr(file_id: int, path: str, media_meta: Optional[str]) -> List[ChunkData]:
    """Image PDF: render pages -> PaddleOCR CPU -> text chunks.
    
    Strategy:
    1. Render each page to image
    2. OCR all pages with PaddleOCR CPU (batch)
    3. Combine OCR text -> unified filter (<500 chars or garbage)
    4. Chunk the combined text
    """
    try:
        import fitz
        from PIL import Image
    except ImportError:
        log.warning("PyMuPDF or Pillow not installed, skipping image PDF")
        return []

    RENDER_DPI = 150  # PDF rendering resolution (150 DPI for OCR quality)

    # Pre-read PDF to memory (avoid repeated I/O for large files)
    try:
        with open(path, 'rb') as f:
            pdf_data = f.read()
        doc = fitz.open(stream=pdf_data, filetype='pdf')
    except Exception as e:
        log.warning(f"Failed to open PDF {path}: {e}")
        return []
    
    total_pages = len(doc)

    import tempfile
    tmp_dir = os.path.join(tempfile.gettempdir(), f'rag_pdf_{file_id}')
    os.makedirs(tmp_dir, exist_ok=True)

    # Render pages (single-threaded, PyMuPDF Document is not thread-safe)
    image_paths = []
    for page_num in range(total_pages):
        try:
            page = doc[page_num]
            pix = page.get_pixmap(dpi=RENDER_DPI)
            # Use JPEG instead of PNG (faster encoding, smaller files)
            tmp_path = os.path.join(tmp_dir, f'page_{page_num+1}.jpg')
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img.save(tmp_path, 'JPEG', quality=85)
            image_paths.append(tmp_path)
        except Exception as e:
            log.warning(f"Error rendering page {page_num+1}: {e}")
        if (page_num + 1) % 100 == 0:
            log.info(f"Rendered {page_num+1}/{total_pages} pages")

    # Close PDF document and free memory
    try:
        doc.close()
    except Exception:
        pass
    del pdf_data

    # Close PDF document
    try:
        doc.close()
    except Exception as e:
        log.warning(f"Error closing PDF: {e}")

    if not image_paths:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        return []

    # Batch OCR all pages (100 pages per batch to avoid timeout)
    BATCH_SIZE = 100
    ocr_results = []
    for batch_start in range(0, len(image_paths), BATCH_SIZE):
        batch_end = min(batch_start + BATCH_SIZE, len(image_paths))
        batch_paths = image_paths[batch_start:batch_end]
        log.info(f"OCR pages {batch_start+1}-{batch_end}/{len(image_paths)}")
        batch_results = _ocr_images_batch(batch_paths)
        ocr_results.extend(batch_results)
    
    # Clean up temp images
    for tmp_path in image_paths:
        try:
            os.remove(tmp_path)
        except OSError:
            pass
    
    # Collect OCR text from all pages
    all_ocr_text = []
    for i, ocr_text in enumerate(ocr_results):
        if ocr_text:
            all_ocr_text.append(f"[Page {i+1}]\n{ocr_text}")

    if not all_ocr_text:
        log.debug(f"No OCR text from PDF: {path}")
        _mark_ocr_insufficient(file_id, path, 'pdf_ocr_empty')
        return []

    combined_text = "\n\n".join(all_ocr_text)
    total_chars = len(combined_text.strip())
    
    # P3: Use lower threshold for PDFs (charts, diagrams, scanned docs)
    threshold = MIN_TEXT_CHARS_LOW
    if total_chars < threshold:
        log.debug(f"PDF OCR too short ({total_chars} chars < {threshold}): {path}")
        _mark_ocr_insufficient(file_id, path, f'pdf_ocr_short_{total_chars}chars')
        return []
    
    if _is_ocr_garbage(combined_text.strip(), min_chars=threshold):
        log.debug(f"PDF OCR garbled/low quality: {path}")
        _mark_ocr_insufficient(file_id, path, 'pdf_ocr_garbled')
        return []

    # Chunk — with per-chunk page range estimation
    short_ctx = _build_context_text(path, media_meta, short=True)
    full_text = f"{short_ctx}\n\nOCR extracted text:\n\n{combined_text}"
    ctx_len = len(full_text) - len(combined_text) - 4  # length of prefix
    total_pages = len(ocr_results)
    
    chunks = chunk_greedy_semantic(full_text, ext='.txt')
    # 保存原始 chunk 长度（不含 overlap），用于页码估算
    raw_chunk_lengths = [len(c) for c in chunks]
    result = []
    for i, c in enumerate(chunks):
        # Estimate page range: map chunk position within combined text
        # 使用原始 chunk 长度（不含 overlap）计算偏移量
        chunk_offset = sum(raw_chunk_lengths[:i]) - ctx_len
        if chunk_offset < 0:
            chunk_offset = 0
        if total_pages > 0:
            start_page = int(chunk_offset / max(len(combined_text), 1) * total_pages) + 1
            end_page = int((chunk_offset + raw_chunk_lengths[i]) / max(len(combined_text), 1) * total_pages) + 1
            start_page = max(1, min(start_page, total_pages))
            end_page = max(start_page, min(end_page, total_pages))
            if start_page == end_page:
                page_str = f"p.{start_page}"
            else:
                page_str = f"p.{start_page}-{end_page}"
        else:
            page_str = None
        
        meta = {'source': 'pdf_ocr', 'ocr_pages': total_pages}
        if page_str:
            meta['pages'] = page_str
        result.append(ChunkData(file_id=file_id, chunk_index=i, text=c, char_count=len(c), meta=meta))
    return result


def _extract_djvu(file_id: int, path: str, media_meta: Optional[str]) -> List[ChunkData]:
    """djvu: 跟 PDF 相同逻辑 — 文本 → OCR → 视觉模型
    
    流程：
    1. djvutxt 提取文本
    2. 文本够长且不是garbage → 直接返回
    3. 文本不够 → 渲染djvu为图片 → PaddleOCR
    4. OCR够长且不是garbage → 返回
    5. OCR不够 → 标记rag_ocr_status → 返回空（等待视觉模型处理）
    """
    # Step 1: 先用 djvutxt 提取文本
    text = ""
    try:
        r = subprocess.run(['djvutxt', path], capture_output=True, text=True, timeout=30)
        if r.returncode == 0 and r.stdout.strip():
            text = r.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    
    # Step 2: 文本够长 → 直接返回
    if text and len(text.strip()) >= MIN_TEXT_CHARS:
        if not _is_ocr_garbage(text.strip(), min_chars=MIN_TEXT_CHARS):
            chunks = chunk_greedy_semantic(text, ext='.txt')
            result = []
            short_ctx = _build_context_text(path, media_meta, short=True)
            for i, chunk_text in enumerate(chunks):
                full_chunk = f"{short_ctx}\n\n{chunk_text}"
                result.append(ChunkData(
                    file_id=file_id, chunk_index=i,
                    text=full_chunk, char_count=len(full_chunk),
                    meta={'source': 'djvu_text'}
                ))
            return result
    
    # Step 3: 文本不够，渲染djvu为图片 → PaddleOCR (via subprocess)
    
    # 获取页数
    try:
        r = subprocess.run(['djvused', '-e', 'n', path], capture_output=True, text=True, timeout=30)
        if r.returncode != 0:
            _mark_ocr_insufficient(file_id, path, 'djvu_info_failed')
            return []
        total_pages = int(r.stdout.strip())
    except (ValueError, subprocess.TimeoutExpired, FileNotFoundError):
        _mark_ocr_insufficient(file_id, path, 'djvu_info_failed')
        return []
    
    if total_pages == 0:
        _mark_ocr_insufficient(file_id, path, 'djvu_zero_pages')
        return []
    
    # 限制最大页数（避免超大文件跑数小时）
    if total_pages > MAX_DJVU_PAGES:
        log.warning(f"djvu has {total_pages} pages, clamping to {MAX_DJVU_PAGES}: {path}")
        total_pages = MAX_DJVU_PAGES
    
    # 渲染全部页面为图片
    import tempfile
    tmp_dir = tempfile.mkdtemp(prefix='rag_djvu_')
    image_paths = []
    
    for page_num in range(1, total_pages + 1):
        out_path = os.path.join(tmp_dir, f'page_{page_num}.png')
        try:
            r = subprocess.run(
                ['ddjvu', f'-page={page_num}', '-format=png', '-resolution=150', path, out_path],
                capture_output=True, timeout=30
            )
            if r.returncode == 0 and os.path.exists(out_path):
                image_paths.append(out_path)
        except subprocess.TimeoutExpired:
            log.warning(f"djvu page {page_num} render timeout: {path}")
    
    if not image_paths:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        _mark_ocr_insufficient(file_id, path, 'djvu_render_failed')
        return []
    
    # PaddleOCR
    ocr_results = _ocr_images_batch(image_paths)
    
    # 清理临时文件
    shutil.rmtree(tmp_dir, ignore_errors=True)
    
    # 收集OCR文本
    all_ocr_text = []
    for i, ocr_text in enumerate(ocr_results):
        if ocr_text:
            all_ocr_text.append(f"[Page {i+1}]\n{ocr_text}")
    
    if not all_ocr_text:
        _mark_ocr_insufficient(file_id, path, 'djvu_ocr_empty')
        return []
    
    combined_text = "\n\n".join(all_ocr_text)
    total_chars = len(combined_text.strip())
    
    # Step 4: OCR够长 → 返回
    if total_chars >= MIN_TEXT_CHARS_LOW and not _is_ocr_garbage(combined_text.strip(), min_chars=MIN_TEXT_CHARS_LOW):
        short_ctx = _build_context_text(path, media_meta, short=True)
        full_text = f"{short_ctx}\n\nOCR extracted text:\n\n{combined_text}"
        ctx_len = len(full_text) - len(combined_text) - 4
        
        chunks = chunk_greedy_semantic(full_text, ext='.txt')
        # 保存原始 chunk 长度（不含 overlap），用于页码估算
        raw_chunk_lengths = [len(c) for c in chunks]
        result = []
        for i, c in enumerate(chunks):
            # 使用原始 chunk 长度（不含 overlap）计算偏移量
            chunk_offset = sum(raw_chunk_lengths[:i]) - ctx_len
            if chunk_offset < 0:
                chunk_offset = 0
            if total_pages > 0:
                start_page = int(chunk_offset / max(len(combined_text), 1) * total_pages) + 1
                end_page = int((chunk_offset + raw_chunk_lengths[i]) / max(len(combined_text), 1) * total_pages) + 1
                start_page = max(1, min(start_page, total_pages))
                end_page = max(start_page, min(end_page, total_pages))
                page_str = f"p.{start_page}" if start_page == end_page else f"p.{start_page}-{end_page}"
            else:
                page_str = None
            
            meta = {'source': 'djvu_ocr', 'ocr_pages': total_pages}
            if page_str:
                meta['pages'] = page_str
            result.append(ChunkData(file_id=file_id, chunk_index=i, text=c, char_count=len(c), meta=meta))
        return result
    
    # Step 5: OCR不够 → 标记，等待视觉模型
    _mark_ocr_insufficient(file_id, path, f'djvu_ocr_short_{total_chars}chars')
    return []


def _extract_image_ocr(file_id: int, path: str, media_meta: str = None) -> List[ChunkData]:
    """Image: PaddleOCR CPU -> text chunks.

    Filter: Use lower threshold for images/screenshots/logos. If insufficient,
    mark in DB as OCR-insufficient and skip (no VLM fallback).
    """
    ext = os.path.splitext(path)[1].lower()
    threshold = MIN_TEXT_CHARS_LOW if ext in LOW_OCR_THRESHOLD_EXTS else MIN_TEXT_CHARS

    ocr_text = _ocr_image(path)

    if not ocr_text:
        log.debug(f"Image OCR empty: {path}")
        _mark_ocr_insufficient(file_id, path, 'ocr_empty')
        return []

    # Filter: insufficient text
    if len(ocr_text.strip()) < threshold:
        log.debug(f"Image OCR too short ({len(ocr_text.strip())} chars < {threshold}): {path}")
        _mark_ocr_insufficient(file_id, path, f'insufficient_text_{len(ocr_text.strip())}chars')
        return []

    if _is_ocr_garbage(ocr_text.strip(), min_chars=threshold):
        log.debug(f"Image OCR garbled/low quality: {path}")
        _mark_ocr_insufficient(file_id, path, 'garbled_ocr')
        return []

    # OCR produced useful text -> chunk it
    ctx = _build_context_text(path, media_meta, short=True)
    full_text = f"{ctx}\n\nOCR extracted text:\n\n{ocr_text}"

    chunks = chunk_greedy_semantic(full_text, ext='.txt')


    return [ChunkData(file_id=file_id, chunk_index=i, text=c, char_count=len(c)) for i, c in enumerate(chunks)]


def _get_audio_duration(path: str) -> float:
    """Get audio duration in seconds using ffprobe."""
    try:
        import subprocess
        r = subprocess.run(
            ['ffprobe', '-v', 'error', '-show_entries', 'format=duration', '-of', 'default=noprint_wrappers=1:nokey=1', path],
            capture_output=True, text=True, timeout=10
        )
        if r.returncode == 0 and r.stdout.strip():
            return float(r.stdout.strip())
    except Exception:
        pass
    return 0

def _extract_audio(file_id: int, path: str, file_size: int = 0, media_meta: str = None) -> List[ChunkData]:
    """Audio: Whisper transcript -> text chunks."""
    result = []
    context_text = _build_context_text(path, media_meta)

    if _is_likely_music_only(path, file_size, media_meta):
        log.debug(f"Skipping Whisper (music-only): {path}")
        result.append(ChunkData(
            file_id=file_id, chunk_index=0,
            text=_build_context_text(path, media_meta),  # just context, no transcript
            char_count=0,
            meta={'source': 'audio_context', 'skip_reason': 'music_only'}
        ))
        return result

    # Pre-check duration with ffprobe to skip very long audio early
    duration_sec = _get_audio_duration(path)
    if duration_sec > 600:  # 10 minutes
        log.warning(f"Audio too long ({duration_sec/60:.1f} min > 10 min), skip transcription: {path}")
        result.append(ChunkData(file_id=file_id, chunk_index=0, text=_build_context_text(path, media_meta), char_count=0, meta={'source': 'audio_context', 'skip_reason': 'too_long'}))
        return result

    cpp_result = _transcribe_with_whisper_cpp(path)
    if cpp_result:
        text, lang = cpp_result
        if len(text.strip()) > 20:
            chunks = chunk_greedy_semantic(text, ext='.txt')
            short_ctx = _build_context_text(path, media_meta, short=True)
            for i, chunk_text in enumerate(chunks):
                full_chunk = f"{short_ctx}\n\n{chunk_text}"
                result.append(ChunkData(
                    file_id=file_id, chunk_index=i,
                    text=full_chunk, char_count=len(full_chunk),
                    meta={'source': 'audio', 'language': lang, 'engine': 'whisper.cpp'}
                ))
            return result

    try:
        # Run faster_whisper in a subprocess to isolate segfaults
        # faster_whisper's PyAV decoder can segfault on corrupted audio files,
        # which would crash the entire worker process and cascade to crash dumps.
        import multiprocessing
        transcribe_ctx = multiprocessing.get_context('fork')
        result_queue = transcribe_ctx.Queue()

        def _whisper_subprocess(q, audio_path):
            """Run transcription in isolated subprocess.
            
            设计权衡：每次子进程启动都重新加载模型（约2-3秒开销），
            但这是为了隔离 segfault 风险（faster_whisper 的 PyAV 解码器
            在损坏的音频文件上会 segfault，会崩溃整个 worker 进程）。
            模型加载开销 vs 稳定性，选择稳定性。
            """
            try:
                from faster_whisper import WhisperModel
                model = WhisperModel("small", device="cpu", compute_type="int8")
                segments, info = model.transcribe(audio_path, beam_size=1)
                duration_sec = info.duration if hasattr(info, 'duration') else 0
                collected = []
                for seg in segments:
                    collected.append(seg.text)
                text = " ".join(collected)
                q.put({'ok': True, 'text': text, 'language': info.language, 'duration': duration_sec})
            except Exception as e:
                q.put({'ok': False, 'error': str(e)})

        proc = transcribe_ctx.Process(target=_whisper_subprocess, args=(result_queue, path))
        proc.start()
        proc.join(timeout=300)  # 5 min timeout

        if proc.is_alive():
            proc.kill()
            proc.join()
            log.warning(f"Whisper subprocess timed out (5 min) for {path}")
        elif proc.exitcode == 0:
            try:
                res = result_queue.get(timeout=1)
            except queue.Empty:
                res = None
            if res and res.get('ok'):
                text = res['text']
                duration_sec = res.get('duration', 0)
                lang = res.get('language', 'unknown')
                log.info(f"Whisper: {duration_sec/60:.1f} min, lang={lang} for {path}")

                if duration_sec > 600:
                    log.warning(f"Audio too long ({duration_sec/60:.1f} min), skip: {path}")
                    result.append(ChunkData(file_id=file_id, chunk_index=0, text=_build_context_text(path, media_meta), char_count=0, meta={'source': 'audio_context', 'skip_reason': 'too_long'}))
                    return result

                if len(text.strip()) < MIN_TEXT_CHARS:
                    log.debug(f"Audio transcript too short ({len(text.strip())} chars): {path}")
                elif not _is_ocr_garbage(text.strip(), min_chars=MIN_TEXT_CHARS):
                    chunks = chunk_greedy_semantic(text, ext='.txt')
                    short_ctx = _build_context_text(path, media_meta, short=True)
                    for i, chunk_text in enumerate(chunks):
                        full_chunk = f"{short_ctx}\n\n{chunk_text}"
                        result.append(ChunkData(
                            file_id=file_id, chunk_index=i,
                            text=full_chunk, char_count=len(full_chunk),
                            meta={'source': 'audio', 'language': lang, 'engine': 'faster-whisper-subprocess'}
                        ))
                    return result
                else:
                    log.debug(f"Audio transcript garbled: {path}")
            else:
                log.debug(f"Whisper subprocess error: {res.get('error', 'unknown')} for {path}")
        elif proc.exitcode != 0:
            log.warning(f"Whisper subprocess crashed (exit={proc.exitcode}) for {path} — skipping")
    except Exception as e:
        log.debug(f"Whisper dispatch failed for {path}: {e}")

    result.append(ChunkData(
        file_id=file_id, chunk_index=0,
        text=_build_context_text(path, media_meta),  # just context, no transcript
        char_count=0,
        meta={'source': 'audio_context'}
    ))
    return result


def _extract_ebook(file_id: int, path: str, ext: str, media_meta: str = None) -> List[ChunkData]:
    """Ebook: extract text + embedded images OCR.
    
    Strategy:
    1. Extract text from ebook format
    2. If text is sufficient -> chunk and return
    3. If text is insufficient -> extract embedded images and OCR
    4. If OCR yields text -> chunk and return
    5. If all fails -> mark as insufficient (same as PDF)
    """
    import tempfile
    tmp_dir = os.path.join(tempfile.gettempdir(), f'rag_ebook_{file_id}')
    os.makedirs(tmp_dir, exist_ok=True)
    
    text = ""
    image_paths = []

    try:
        if ext.lower() == '.epub':
            text, image_paths = _extract_epub_content(path, tmp_dir)

        elif ext.lower() == '.mobi':
            text, image_paths = _extract_mobi_content(path, tmp_dir)
            
        elif ext.lower() == '.azw3':
            text, image_paths = _extract_azw3_content(path, tmp_dir)
            
        elif ext.lower() == '.fb2':
            text, image_paths = _extract_fb2_content(path, tmp_dir)
            
        else:
            text = _read_file(path) or ""

    except Exception as e:
        log.warning(f"Ebook extraction error for {path}: {e}")
        text = _read_file(path) or ""

    # P1: Check if text extraction was sufficient
    if len(text.strip()) >= MIN_TEXT_CHARS and not _is_ocr_garbage(text.strip(), min_chars=MIN_TEXT_CHARS):
        # Clean up temp dir
        shutil.rmtree(tmp_dir, ignore_errors=True)
        chunks = chunk_greedy_semantic(text, ext='.txt')
        result = []
        short_ctx = _build_context_text(path, media_meta, short=True)
        for i, chunk_text in enumerate(chunks):
            full_chunk = f"{short_ctx}\n\n{chunk_text}"
            result.append(ChunkData(
                file_id=file_id, chunk_index=i,
                text=full_chunk, char_count=len(full_chunk),
                meta={'source': 'ebook'}
            ))
        return result

    # P2: Text insufficient, try OCR on embedded images
    if image_paths:
        log.info(f"Ebook text insufficient ({len(text.strip())} chars), OCRing {len(image_paths)} images: {path}")
        
        # Batch OCR (100 images per batch)
        BATCH_SIZE = 100
        ocr_results = []
        for batch_start in range(0, len(image_paths), BATCH_SIZE):
            batch_end = min(batch_start + BATCH_SIZE, len(image_paths))
            batch_paths = image_paths[batch_start:batch_end]
            log.info(f"Ebook OCR pages {batch_start+1}-{batch_end}/{len(image_paths)}")
            batch_results = _ocr_images_batch(batch_paths)
            ocr_results.extend(batch_results)
        
        # Collect OCR text
        all_ocr_text = []
        for i, ocr_text in enumerate(ocr_results):
            if ocr_text:
                all_ocr_text.append(f"[Page {i+1}]\n{ocr_text}")
        
        combined_ocr = "\n\n".join(all_ocr_text)
        
        # Clean up temp images
        for tmp_path in image_paths:
            try:
                os.remove(tmp_path)
            except OSError:
                pass
        shutil.rmtree(tmp_dir, ignore_errors=True)
        
        # Check OCR results
        if not combined_ocr.strip():
            _mark_ocr_insufficient(file_id, path, 'ebook_ocr_empty')
            return []
        
        if len(combined_ocr.strip()) < MIN_TEXT_CHARS:
            _mark_ocr_insufficient(file_id, path, f'ebook_ocr_short_{len(combined_ocr.strip())}chars')
            return []
        
        if _is_ocr_garbage(combined_ocr.strip(), min_chars=MIN_TEXT_CHARS):
            _mark_ocr_insufficient(file_id, path, 'ebook_ocr_garbled')
            return []
        
        # OCR succeeded, chunk the text
        chunks = chunk_greedy_semantic(combined_ocr, ext='.txt')
        result = []
        short_ctx = _build_context_text(path, media_meta, short=True)
        for i, chunk_text in enumerate(chunks):
            full_chunk = f"{short_ctx}\n\n{chunk_text}"
            result.append(ChunkData(
                file_id=file_id, chunk_index=i,
                text=full_chunk, char_count=len(full_chunk),
                meta={'source': 'ebook_ocr'}
            ))
        return result
    
    # P3: No images to OCR, text was insufficient
    # Clean up
    shutil.rmtree(tmp_dir, ignore_errors=True)
    
    if len(text.strip()) < MIN_TEXT_CHARS:
        _mark_ocr_insufficient(file_id, path, f'ebook_short_{len(text.strip())}chars')
        return []
    
    if _is_ocr_garbage(text.strip(), min_chars=MIN_TEXT_CHARS):
        _mark_ocr_insufficient(file_id, path, 'ebook_garbled')
        return []
    
    # Should not reach here, but just in case
    return []

def _extract_epub_content(epub_path: str, tmp_dir: str) -> tuple:
    """Extract text and images from EPUB.
    
    Returns: (text, image_paths)
    """
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    import uuid
    
    text = ""
    image_paths = []
    epub_tmp = None
    
    try:
        # ebooklib 不支持非 ASCII 路径
        epub_tmp = os.path.join(tmp_dir, f'temp_{uuid.uuid4().hex}.epub')
        shutil.copyfile(epub_path, epub_tmp)
        
        book = epub.read_epub(epub_tmp)
        
        # Extract text from HTML documents
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            content_bytes = item.get_content().decode('utf-8', errors='replace')
            if not content_bytes.strip():
                continue
            
            # Remove style/script but keep img tags for counting
            content_bytes = re.sub(r'<style[^>]*>.*?</style>', '', content_bytes, flags=re.DOTALL | re.IGNORECASE)
            content_bytes = re.sub(r'<script[^>]*>.*?</script>', '', content_bytes, flags=re.DOTALL | re.IGNORECASE)
            content_bytes = re.sub(r'<link[^>]*/?>', '', content_bytes, flags=re.IGNORECASE)
            
            soup = BeautifulSoup(content_bytes, 'html.parser')
            
            # Count images in this document
            img_count = len(soup.find_all('img'))
            
            # Extract text
            for tag in soup.find_all(['h1', 'h2', 'h3', 'h4']):
                tag.replace_with(f"\n\n{tag.get_text()}\n\n")
            for tag in soup.find_all('li'):
                tag.replace_with(f"\n- {tag.get_text()}")
            clean = soup.get_text(separator='\n')
            clean = re.sub(r'\n{3,}', '\n\n', clean)
            text += clean + "\n"
        
        # If text is insufficient, extract images for OCR
        if len(text.strip()) < MIN_TEXT_CHARS:
            # Extract all images
            for item in book.get_items_of_type(ebooklib.ITEM_IMAGE):
                img_data = item.get_content()
                img_name = item.get_name()
                # Save to temp dir
                img_path = os.path.join(tmp_dir, os.path.basename(img_name))
                with open(img_path, 'wb') as f:
                    f.write(img_data)
                image_paths.append(img_path)
    
    except Exception as e:
        log.warning(f"EPUB extraction error: {e}")
        text = _read_file(epub_path) or ""
    
    return text, image_paths


def _extract_mobi_content(mobi_path: str, tmp_dir: str) -> tuple:
    """Extract text and images from MOBI.
    
    MOBI structure:
    - mobi7/book.html: fallback HTML (may be image-heavy)
    - mobi8/: KF8 format with part*.xhtml and images
    
    Returns: (text, image_paths)
    """
    import mobi
    from bs4 import BeautifulSoup
    
    text = ""
    image_paths = []
    tempdir = None
    
    try:
        tempdir, filepath = mobi.extract(mobi_path)
        
        # P1: Try mobi8 (KF8) first - has better structure
        mobi8_dir = os.path.join(tempdir, 'mobi8')
        if os.path.isdir(mobi8_dir):
            # Extract text from all xhtml files
            xhtml_files = []
            for root, dirs, files in os.walk(mobi8_dir):
                for f in files:
                    if f.endswith(('.xhtml', '.html', '.htm')):
                        xhtml_files.append(os.path.join(root, f))
            
            # Sort by name to maintain order
            xhtml_files.sort()
            
            for xhtml_path in xhtml_files:
                try:
                    with open(xhtml_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                    soup = BeautifulSoup(content, 'html.parser')
                    page_text = soup.get_text(separator='\n').strip()
                    if page_text:
                        text += page_text + "\n"
                except Exception:
                    pass
        
        # P2: Fallback to mobi7 if mobi8 text is insufficient
        if len(text.strip()) < MIN_TEXT_CHARS:
            mobi7_dir = os.path.join(tempdir, 'mobi7')
            if os.path.isdir(mobi7_dir):
                for f in os.listdir(mobi7_dir):
                    if f.endswith(('.html', '.htm')):
                        html_path = os.path.join(mobi7_dir, f)
                        try:
                            with open(html_path, 'r', encoding='utf-8', errors='replace') as fh:
                                soup = BeautifulSoup(fh.read(), 'html.parser')
                                text += soup.get_text(separator='\n') + "\n"
                        except Exception:
                            pass
        
        # P3: If text still insufficient, extract images for OCR
        if len(text.strip()) < MIN_TEXT_CHARS:
            # Find all images in extracted directory
            for root, dirs, files in os.walk(tempdir):
                for f in files:
                    if f.lower().endswith(('.jpg', '.jpeg', '.png', '.gif')):
                        src_path = os.path.join(root, f)
                        # Copy to tmp_dir to avoid cleanup issues
                        dst_path = os.path.join(tmp_dir, f)
                        try:
                            shutil.copy2(src_path, dst_path)
                            image_paths.append(dst_path)
                        except Exception:
                            pass
            
            # Sort images by name
            image_paths.sort()
    
    except Exception as e:
        log.warning(f"MOBI extraction error: {e}")
        text = _read_file(mobi_path) or ""
    
    finally:
        # Clean up mobi extract directory
        if tempdir and os.path.isdir(tempdir):
            shutil.rmtree(tempdir, ignore_errors=True)
    
    return text, image_paths


def _extract_azw3_content(azw3_path: str, tmp_dir: str) -> tuple:
    """Extract text and images from AZW3.
    
    Uses ebook-convert if available, otherwise falls back to mobi extraction.
    
    Returns: (text, image_paths)
    """
    import subprocess
    
    text = ""
    image_paths = []
    
    # Try ebook-convert first
    out_txt = os.path.join(tmp_dir, 'converted.txt')
    try:
        r = subprocess.run(['ebook-convert', azw3_path, out_txt],
                          capture_output=True, text=True, timeout=120)
        if r.returncode == 0 and os.path.exists(out_txt):
            with open(out_txt, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    
    # Fallback: try mobi extraction (azw3 is similar to mobi)
    if len(text.strip()) < MIN_TEXT_CHARS:
        try:
            import mobi
            tempdir, filepath = mobi.extract(azw3_path)
            
            # Extract text from HTML files
            for root, dirs, files in os.walk(tempdir):
                for f in files:
                    if f.endswith(('.html', '.htm', '.xhtml')):
                        html_path = os.path.join(root, f)
                        try:
                            from bs4 import BeautifulSoup
                            with open(html_path, 'r', encoding='utf-8', errors='replace') as fh:
                                soup = BeautifulSoup(fh.read(), 'html.parser')
                                text += soup.get_text(separator='\n') + "\n"
                        except Exception:
                            pass
            
            # If text still insufficient, extract images
            if len(text.strip()) < MIN_TEXT_CHARS:
                for root, dirs, files in os.walk(tempdir):
                    for f in files:
                        if f.lower().endswith(('.jpg', '.jpeg', '.png', '.gif')):
                            src_path = os.path.join(root, f)
                            dst_path = os.path.join(tmp_dir, f)
                            try:
                                shutil.copy2(src_path, dst_path)
                                image_paths.append(dst_path)
                            except Exception:
                                pass
                image_paths.sort()
            
            shutil.rmtree(tempdir, ignore_errors=True)
        except Exception as e:
            log.warning(f"AZW3 extraction error: {e}")
    
    return text, image_paths


def _extract_fb2_content(fb2_path: str, tmp_dir: str) -> tuple:
    """Extract text and images from FB2 (FictionBook).
    
    FB2 is XML-based, images are base64 encoded inline.
    
    Returns: (text, image_paths)
    """
    from bs4 import BeautifulSoup
    import base64
    
    text = ""
    image_paths = []
    
    try:
        with open(fb2_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        
        soup = BeautifulSoup(content, 'xml')
        
        # Extract text from body sections
        body = soup.find('body')
        if body:
            text = body.get_text(separator='\n')
        else:
            # Fallback: extract all text
            text = soup.get_text(separator='\n')
        
        # If text insufficient, extract base64 images
        if len(text.strip()) < MIN_TEXT_CHARS:
            images = soup.find_all('image')
            for i, img in enumerate(images):
                # FB2 images are referenced by id, actual data in <binary>
                img_id = img.get('l:href') or img.get('{http://www.w3.org/1999/xlink}href')
                if img_id and img_id.startswith('#'):
                    img_id = img_id[1:]
                
                # Find the binary element
                binary = soup.find('binary', {'id': img_id})
                if binary and binary.string:
                    try:
                        # Decode base64
                        img_data = base64.b64decode(binary.string)
                        # Determine extension from content-type
                        content_type = binary.get('content-type', 'image/jpeg')
                        ext = '.jpg' if 'jpeg' in content_type else '.png'
                        img_path = os.path.join(tmp_dir, f'fb2_img_{i}{ext}')
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        image_paths.append(img_path)
                    except Exception:
                        pass
    
    except Exception as e:
        log.warning(f"FB2 extraction error: {e}")
        text = _read_file(fb2_path) or ""
    
    return text, image_paths

def _ocr_embedded_office_images(zip_path: str, media_prefix: str, results: list):
    """P2: Extract embedded images from Office file and OCR them.
    
    Args:
        zip_path: Path to the Office file (docx/pptx are zip files)
        media_prefix: Subdirectory prefix (e.g., 'word/media/', 'ppt/media/')
        results: List to append OCR results to
    
    Size filtering (same as standalone image files):
        - Skip < MIN_MEDIA_BYTES (10KB): too small, likely icons/decorations
        - Skip > MAX_IMAGE_BYTES (100MB): too large, OCR timeout risk
    """
    tmp_dir = None
    try:
        import zipfile
        with zipfile.ZipFile(zip_path, 'r') as zf:
            img_names = [n for n in zf.namelist()
                         if n.startswith(media_prefix)
                         and n.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff', '.tif'))]
            if img_names:
                tmp_dir = tempfile.mkdtemp(prefix='rag_office_img_')
                img_paths = []
                skipped_small = 0
                skipped_large = 0
                for img_name in img_names:
                    info = zf.getinfo(img_name)
                    # Size filtering: same rules as standalone image files
                    if info.file_size < MIN_MEDIA_BYTES:
                        skipped_small += 1
                        continue
                    if info.file_size > MAX_IMAGE_BYTES:
                        skipped_large += 1
                        continue
                    zf.extract(img_name, tmp_dir)
                    img_paths.append(os.path.join(tmp_dir, img_name))
                if skipped_small or skipped_large:
                    log.debug(f"Office embedded images filtered: {skipped_small} too small (<{MIN_MEDIA_BYTES}B), {skipped_large} too large (>{MAX_IMAGE_BYTES//1024//1024}MB)")
                if img_paths:
                    ocr_results = _ocr_images_batch(img_paths)
                    for i, ocr_t in enumerate(ocr_results):
                        if ocr_t and len(ocr_t.strip()) >= MIN_TEXT_CHARS_LOW:
                            results.append(f"[Embedded Image {i+1}: {os.path.basename(img_paths[i])}]\n{ocr_t}")
    except Exception as e:
        log.debug(f"Office embedded image OCR failed: {e}")
    finally:
        if tmp_dir:
            shutil.rmtree(tmp_dir, ignore_errors=True)


def _extract_office(file_id: int, path: str, ext: str, media_meta: str = None) -> List[ChunkData]:
    """Office documents: extract text -> greedy chunk.
    
    P2: Also extracts embedded images from docx/pptx and OCRs them.
    """
    text = ""
    ext = ext.lower()
    embedded_ocr_texts = []  # P2: Collect OCR text from embedded images

    if ext == '.docx':
        try:
            from docx import Document
            from docx.oxml.ns import qn
            doc = Document(path)
            body = doc.element.body
            parts = []
            para_idx = 0
            for child in body.iterchildren():
                if child.tag == qn('w:tbl'):
                    for row_elem in child.iterchildren(qn('w:tr')):
                        cells = []
                        for cell_elem in row_elem.iterchildren(qn('w:tc')):
                            cell_texts = []
                            for p in cell_elem.iterchildren(qn('w:p')):
                                for r in p.iterchildren(qn('w:r')):
                                    for t in r.iterchildren(qn('w:t')):
                                        if t.text:
                                            cell_texts.append(t.text)
                            cells.append(''.join(cell_texts))
                        if cells:
                            parts.append(' | '.join(cells))
                    parts.append('')
                elif child.tag == qn('w:p'):
                    if para_idx < len(doc.paragraphs):
                        parts.append(doc.paragraphs[para_idx].text)
                        para_idx += 1
            text = '\n'.join(parts) if parts else '\n'.join(p.text for p in doc.paragraphs)
            
            # P2: Extract embedded images from docx (word/media/*.png etc.) and OCR them
            _ocr_embedded_office_images(path, 'word/media/', embedded_ocr_texts)
        except ImportError:
            text = _read_file(path) or ""
        except Exception as e:
            log.debug(f"docx extract error: {e}")
            text = _read_file(path) or ""
    elif ext == '.doc':
        try:
            r = subprocess.run(['antiword', path], capture_output=True, text=True, timeout=10)
            if r.returncode == 0 and r.stdout.strip():
                text = r.stdout
            else:
                text = _read_file(path) or ""
        except FileNotFoundError:
            text = _read_file(path) or ""
        except subprocess.TimeoutExpired:
            text = _read_file(path) or ""
    elif ext in ('.xlsx', '.xls'):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True)
            for ws in wb.worksheets:
                text += f"\n--- Sheet: {ws.title} ---\n"
                for row in ws.iter_rows(values_only=True):
                    text += " | ".join(str(c) for c in row if c is not None) + "\n"
            wb.close()
        except Exception:
            if ext == '.xls':
                try:
                    import xlrd
                    wb = xlrd.open_workbook(path)
                    for sheet in wb.sheets():
                        text += f"\n--- Sheet: {sheet.name} ---\n"
                        for row_idx in range(sheet.nrows):

                            text += " | ".join(str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)) + "\n"
                except ImportError:
                    pass
    elif ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(path)
            for idx, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {idx} ---\n"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text_frame.text + "\n"
            
            # P2: Extract embedded images from pptx (ppt/media/*.png etc.) and OCR them
            _ocr_embedded_office_images(path, 'ppt/media/', embedded_ocr_texts)
        except ImportError:
            # Fallback: try libreoffice or catppt
            try:
                r = subprocess.run(['catppt', '-d', 'utf-8', path], capture_output=True, text=True, timeout=30)
                if r.returncode == 0 and r.stdout.strip():
                    text = r.stdout
                else:
                    text = _read_file(path) or ""
            except Exception:
                text = _read_file(path) or ""

    elif ext == '.ppt':
        import tempfile, glob
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                r = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pptx', '--outdir', tmpdir, path],
                    capture_output=True, text=True, timeout=60
                )
                if r.returncode == 0:
                    pptx_files = glob.glob(os.path.join(tmpdir, '*.pptx'))
                    if pptx_files:
                        from pptx import Presentation
                        prs = Presentation(pptx_files[0])
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text += shape.text_frame.text + "\n"
            if not text.strip():
                raise ValueError("No text")
        except Exception:
            try:
                r2 = subprocess.run(['catppt', '-d', 'utf-8', path], capture_output=True, text=True, timeout=10)
                if r2.returncode == 0 and r2.stdout.strip():
                    text = r2.stdout
                else:
                    text = _read_file(path) or ""
            except Exception:
                text = _read_file(path) or ""
    elif ext == '.rtf':
        text = _read_file(path) or ""
        text = re.sub(r'\\[a-z]+\d*\s?', '', text)
        text = re.sub(r'[{}]', '', text)
    else:
        text = _read_file(path) or ""

    # P2: Append embedded image OCR text to main text
    if embedded_ocr_texts:
        text += "\n\n=== Embedded Images (OCR) ===\n" + "\n\n".join(embedded_ocr_texts)

    # P3: Use lower threshold for office docs with embedded images (charts, diagrams)
    effective_threshold = MIN_TEXT_CHARS_LOW if embedded_ocr_texts else MIN_TEXT_CHARS
    if len(text.strip()) < effective_threshold:
        _mark_ocr_insufficient(file_id, path, f'office_short_{len(text.strip())}chars')
        return []

    # Garbled text check
    if _is_ocr_garbage(text.strip(), min_chars=effective_threshold):
        _mark_ocr_insufficient(file_id, path, 'office_garbled')
        return []

    ctx = _build_context_text(path, media_meta)
    full_text = ctx + "\n\n" + text
    chunks = chunk_greedy_semantic(full_text, ext='.txt')


    return [ChunkData(file_id=file_id, chunk_index=i, text=c, char_count=len(c)) for i, c in enumerate(chunks)]


def _safe_run_archive(args, timeout=15):
    """subprocess.run for archive tools — uses bytes mode to avoid UTF-8 decode errors.
    
    Windows ZIP files often contain GBK-encoded filenames in their listing output.
    Python's text=True with default UTF-8 encoding will crash on these.
    """
    r = subprocess.run(args, capture_output=True, timeout=timeout)
    stdout = r.stdout.decode('utf-8', errors='replace') if r.stdout else ''
    stderr = r.stderr.decode('utf-8', errors='replace') if r.stderr else ''
    # Simulate CompletedProcess with decoded strings
    r.stdout_text = stdout
    r.stderr_text = stderr
    return r


def _extract_archive(file_id: int, path: str, media_meta: str = None) -> List[ChunkData]:
    """Archive: extract recursive file listing as RAG text.
    
    - Skips password-protected archives (marks in DB)
    - Recursively lists all files including subdirectories
    - Includes file size and extension info
    - Falls back to context-only for unopenable archives
    - Handles non-UTF-8 (GBK, etc.) filenames in archives gracefully
    """
    ctx = _build_context_text(path, media_meta)
    ext = os.path.splitext(path)[1].lower()
    listing = ""
    is_encrypted = False

    try:
        if ext == '.zip':
            # Check for encryption first: unzip -t tests integrity
            r = _safe_run_archive(['unzip', '-t', path])
            stderr_lower = r.stderr_text.lower()
            if 'incorrect password' in stderr_lower or 'cannot find or open' in stderr_lower:
                is_encrypted = True
            elif r.returncode == 82:  # unzip exit code 82 = encrypted
                is_encrypted = True
            else:
                # Recursive listing with sizes: unzip -Z -S (show sizes, dates, full paths)
                r = _safe_run_archive(['unzip', '-Z', '-S', path])
                if r.returncode == 0 and r.stdout_text.strip():
                    listing = r.stdout_text[:40000]
                else:
                    # Fallback: simple listing
                    r2 = _safe_run_archive(['unzip', '-l', path])
                    listing = r2.stdout_text[:40000]

        elif ext in ('.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tar.xz') or path.endswith('.tar.gz'):
            # tar cannot be encrypted, just list
            r = _safe_run_archive(['tar', '-tvf', path])
            listing = r.stdout_text[:40000]

        elif ext == '.7z':
            # Check encryption: 7z l shows "Encrypted = +" in tech mode
            r = _safe_run_archive(['7z', 'l', '-slt', path])
            if 'Encrypted = +' in r.stdout_text:
                is_encrypted = True
            else:
                # Normal listing with recursive subdirs
                r2 = _safe_run_archive(['7z', 'l', '-r', path])
                listing = r2.stdout_text[:40000]

        elif ext == '.rar':
            # Try 7z first, then unrar
            try:
                r = _safe_run_archive(['7z', 'l', '-slt', path])
                if 'Encrypted = +' in r.stdout_text:
                    is_encrypted = True
                elif r.returncode == 0:
                    r2 = _safe_run_archive(['7z', 'l', '-r', path])
                    listing = r2.stdout_text[:40000]
                else:
                    r3 = _safe_run_archive(['unrar', 'l', path])
                    listing = r3.stdout_text[:40000]
            except FileNotFoundError:
                r = _safe_run_archive(['unrar', 'l', path])
                listing = r.stdout_text[:40000]

        elif ext == '.gz' and not path.endswith('.tar.gz'):
            # Single .gz file — try to read content
            import gzip
            try:
                with gzip.open(path, 'rt', errors='replace') as gf:
                    content = gf.read(40000)
                    if content and len(content.strip()) >= MIN_TEXT_CHARS:
                        listing = f"[Gzip content preview]\n{content[:40000]}"
            except Exception:
                listing = f"[Gzip file: {os.path.basename(path)}]"

    except subprocess.TimeoutExpired:
        log.debug(f"Archive listing timeout: {path}")
        listing = "[Archive listing timed out]"
    except (FileNotFoundError, OSError) as e:
        log.debug(f"Archive tool not found: {e}")
        listing = "[Archive tool unavailable]"

    # Handle encrypted archives
    if is_encrypted:
        log.debug(f"Archive encrypted, skipping: {path}")
        _mark_ocr_insufficient(file_id, path, 'archive_encrypted')
        # Still index the filename and basic context
        return [ChunkData(
            file_id=file_id, chunk_index=0,
            text=ctx + "\n\n[Encrypted archive - contents not accessible]\n",
            char_count=len(ctx) + 50,
            meta={'source': 'archive_encrypted'}
        )]

    # Build full text: context + listing
    full_text = ctx + "\n\n=== Archive Contents ===\n" + listing
    if not listing.strip():
        full_text = ctx + "\n\n[Empty or unreadable archive]"

    # Chunk if large
    chunks = chunk_greedy_semantic(full_text, ext='.txt')
    return [ChunkData(
        file_id=file_id, chunk_index=i, text=c, char_count=len(c),
        meta={'source': 'archive'}
    ) for i, c in enumerate(chunks)]


def _extract_text_fallback(file_id: int, path: str, media_meta: str = None) -> List[ChunkData]:
    """Unknown file type: try reading as text."""
    ctx = _build_context_text(path, media_meta)
    text = _read_file(path)
    if not text or len(text.strip()) < MIN_TEXT_CHARS:
        return [ChunkData(
            file_id=file_id, chunk_index=0,
            text=ctx, char_count=len(ctx),
            meta={'source': 'unknown'}
        )]

    # Binary content check: null bytes indicate non-text file
    if '\x00' in text:
        log.debug(f"Skipping binary content (null bytes): {path}")
        return []

    # Garbled text check
    if _is_ocr_garbage(text.strip(), min_chars=MIN_TEXT_CHARS):
        _mark_ocr_insufficient(file_id, path, 'fallback_garbled')
        return []

    full_text = ctx + "\n\n" + text
    chunks = chunk_greedy_semantic(full_text, ext='.txt')
    return [ChunkData(
        file_id=file_id, chunk_index=i, text=c, char_count=len(c),
        meta={'source': 'unknown'}
    ) for i, c in enumerate(chunks)]